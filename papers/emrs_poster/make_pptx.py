#!/usr/bin/env python3
"""Generate EMRS poster as PPTX (A0 portrait, 33.11" × 46.81").

Usage:
    python papers/emrs_poster/make_pptx.py
Output:
    papers/emrs_poster/emrs_poster.pptx
"""

import io, os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))

# ── Slide dimensions (A0 portrait) ────────────────────────
SW = 33.11
SH = 46.81

# ── Colour palette ─────────────────────────────────────────
NAVY      = RGBColor(0x1A, 0x3A, 0x5C)
MIDBLUE   = RGBColor(0x2C, 0x6F, 0xAD)
LTBLUE    = RGBColor(0xD6, 0xE8, 0xF7)
ORANGE    = RGBColor(0xC0, 0x56, 0x2A)
LTORANGE  = RGBColor(0xFA, 0xE5, 0xD9)
LTGRAY    = RGBColor(0xF4, 0xF4, 0xF6)
DKGRAY    = RGBColor(0x3A, 0x3A, 0x3A)
GREEN     = RGBColor(0x2E, 0x7D, 0x32)
LTGREEN   = RGBColor(0xE8, 0xF5, 0xE9)
ROWBLUE   = RGBColor(0xDC, 0xF0, 0xFB)
ROWRED    = RGBColor(0xFD, 0xEC, 0xEA)
ROWYELLOW = RGBColor(0xFF, 0xFD, 0xE7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
RED_TEXT  = RGBColor(0xC0, 0x00, 0x00)

# ── Shorthand ──────────────────────────────────────────────
def I(v): return Inches(v)


# ──────────────────────────────────────────────────────────
# Low-level helpers
# ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, border=None, bw=0.8):
    shape = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(bw)
    else:
        shape.line.fill.background()
    return shape


def add_tb(slide, l, t, w, h, paras,
           ml=0.10, mr=0.10, mt=0.06, mb=0.04):
    """Add textbox. paras = list of dicts: text/size/bold/italic/color/align/space_before."""
    box = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = I(ml); tf.margin_right = I(mr)
    tf.margin_top  = I(mt); tf.margin_bottom = I(mb)
    first = True
    for pg in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = pg.get('align', PP_ALIGN.LEFT)
        sb = pg.get('space_before', 0)
        if sb: p.space_before = Pt(sb)
        txt = pg.get('text', '')
        if txt:
            run = p.add_run()
            run.text = txt
            run.font.size   = Pt(pg.get('size', 14))
            run.font.bold   = pg.get('bold', False)
            run.font.italic = pg.get('italic', False)
            run.font.color.rgb = pg.get('color', DKGRAY)
    return box


def section_box(slide, l, t, w, h,
                title, title_bg, title_fg=WHITE,
                body_bg=LTGRAY, title_h=0.42, title_sz=20):
    """Draw titled box. Returns body_top (inches)."""
    bc = title_bg
    add_rect(slide, l, t, w, h, fill=body_bg, border=bc, bw=1.0)
    add_rect(slide, l, t, w, title_h, fill=title_bg)
    add_tb(slide, l+0.06, t+0.02, w-0.12, title_h-0.02,
           [{'text': title, 'size': title_sz, 'bold': True, 'color': title_fg}],
           ml=0.12, mr=0.06, mt=0.04, mb=0.02)
    return t + title_h   # body_top


def px_img(slide, path, l, t, w):
    """Add picture at width w; returns actual height in inches."""
    pic = slide.shapes.add_picture(path, I(l), I(t), width=I(w))
    return pic.height / 914400.0


# ──────────────────────────────────────────────────────────
# Table helpers
# ──────────────────────────────────────────────────────────

def _cell_bg(cell, color):
    tc = cell._tc
    pr = tc.get_or_add_tcPr()
    for c in pr.findall(qn('a:solidFill')): pr.remove(c)
    sf = etree.SubElement(pr, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr'))
    sc.set('val', str(color))


def _cell_txt(cell, text, size=12, bold=False, italic=False,
              color=DKGRAY, align=PP_ALIGN.CENTER):
    cell.text = ''
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color


def make_survey_table(slide, l, t, w, h):
    nrows, ncols = 9, 5
    tbl = slide.shapes.add_table(nrows, ncols, I(l), I(t), I(w), I(h)).table
    props = [0.33, 0.13, 0.17, 0.17, 0.13]
    s = sum(props)
    for i, p in enumerate(props):
        tbl.columns[i].width = int(I(w) * p / s)

    hdrs = ['Stack', 'CBO (eV)', 'E_res (meV)', 'V_peak (V)', 'PVCR']
    for j, h_ in enumerate(hdrs):
        c = tbl.cell(0, j)
        _cell_bg(c, NAVY)
        _cell_txt(c, h_, size=13, bold=True, color=WHITE)

    rows = [
        ('In₂O₃/Al₂O₃',              '2.8',  '>2000', '4–5',  '~1†',  ROWRED),
        ('In₂O₃/HfO₂',               '2.0',  '>2000', '3–4',  '~1†',  ROWRED),
        ('In₂O₃/β-Ga₂O₃ (2 nm)',     '0.07', '30',    '1.43', '1.33', ROWYELLOW),
        ('In₂O₃/β-Ga₂O₃ (1 nm)',     '0.07', '53',    '1.48', '1.05', ROWYELLOW),
        ('In₂O₃/κ-Ga₂O₃ (2 nm)',     '0.45', '278',   '1.18', '3.55', LTBLUE),
        ('In₂O₃/κ-Ga₂O₃ (1 nm)',     '0.45', '177',   '1.10', '3.50', LTBLUE),
        ('In₂O₃/κ-Ga₂O₃ (asym)',     '0.45', '402',   '1.33', '2.80', LTBLUE),
        ('ZnO/Mg₀.₃Zn₀.₇O (2 nm)',   '0.47', '290',   '1.10', '6.32', ROWBLUE),
    ]
    for i, (*vals, bg) in enumerate(rows):
        is_sel = (i == 7)
        for j, v in enumerate(vals):
            c = tbl.cell(i+1, j)
            _cell_bg(c, bg)
            _cell_txt(c, v, size=12, bold=is_sel,
                      color=NAVY if is_sel else BLACK,
                      align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)


def make_candidate_table(slide, l, t, w, h):
    tbl = slide.shapes.add_table(5, 6, I(l), I(t), I(w), I(h)).table
    props = [0.31, 0.13, 0.13, 0.13, 0.13, 0.15]
    s = sum(props)
    for i, p in enumerate(props):
        tbl.columns[i].width = int(I(w) * p / s)

    hdrs = ['Stack', 'CBO (eV)', 'm*/m₀', 'NDR', 'BEOL', 'Ref.']
    for j, h_ in enumerate(hdrs):
        c = tbl.cell(0, j)
        _cell_bg(c, NAVY)
        _cell_txt(c, h_, size=13, bold=True, color=WHITE)

    rows = [
        ('AlGaAs/GaAs',         '0.23', '0.067', '✓', '✗', '[8]',  WHITE),
        ('GaN/AlN',             '1.8',  '0.20',  '✓', '✗', '[9]',  WHITE),
        ('In₂O₃/κ-Ga₂O₃',      '0.45', '0.30',  '✗', '✓', '[10]', LTBLUE),
        ('ZnO/Mg₀.₃Zn₀.₇O',    '0.47', '0.28',  '✓', '✓', '[3,4]',ROWBLUE),
    ]
    for i, (*vals, bg) in enumerate(rows):
        is_sel = (i == 3)
        for j, v in enumerate(vals):
            c = tbl.cell(i+1, j)
            _cell_bg(c, bg)
            col = (GREEN if v == '✓' else
                   RED_TEXT if v == '✗' else
                   (NAVY if is_sel else BLACK))
            _cell_txt(c, v, size=12, bold=(is_sel or v in '✓✗'),
                      color=col,
                      align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)


def make_ald_table(slide, l, t, w, h):
    tbl = slide.shapes.add_table(5, 4, I(l), I(t), I(w), I(h)).table
    props = [0.22, 0.20, 0.42, 0.16]
    s = sum(props)
    for i, p in enumerate(props):
        tbl.columns[i].width = int(I(w) * p / s)

    hdrs = ['Layer', 'Temp.', 'Precursors', 'Rate']
    for j, h_ in enumerate(hdrs):
        c = tbl.cell(0, j)
        _cell_bg(c, NAVY)
        _cell_txt(c, h_, size=12, bold=True, color=WHITE)

    rows = [
        ('ZnO (well)',           'RT–350 °C',  'DEZ + H₂O',              '1.5–2.0 Å/cyc', WHITE),
        ('MgₓZn₁₋ₓO (barrier)', 'RT–350 °C',  'Mg(Cp)₂ + DEZ + H₂O',   '1.0–1.5 Å/cyc', LTBLUE),
        ('In₂O₃',               '150–275 °C', 'In(DMAMP)₂(OiPr) + H₂O', '0.3–0.5 Å/cyc', LTGRAY),
        ('κ-Ga₂O₃',             '200–350 °C', 'TMGa + O₂ plasma',        '~0.5 Å/cyc',    LTGRAY),
    ]
    for i, (*vals, bg) in enumerate(rows):
        for j, v in enumerate(vals):
            c = tbl.cell(i+1, j)
            _cell_bg(c, bg)
            _cell_txt(c, v, size=11, bold=(i < 2),
                      color=DKGRAY,
                      align=PP_ALIGN.LEFT if j in (0, 2) else PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────
# Design-rule bar chart
# ──────────────────────────────────────────────────────────

def make_bar_chart(out_path):
    fig, ax = plt.subplots(figsize=(8, 4.8))

    ax.axvspan(0.05, 0.40, color='#2E7D32', alpha=0.18)
    ax.axvline(0.05, color='#2E7D32', lw=1.5, ls='--')
    ax.axvline(0.40, color='#2E7D32', lw=1.5, ls='--')
    ax.text(0.225, 5.35, 'Vibration\nwindow\n50–400 mV',
            ha='center', va='top', fontsize=10, color='#1B5E20')

    high_stacks = [
        (3.3, 4.3, 'In₂O₃/Al₂O₃'),
        (4.5, 3.6, 'In₂O₃/HfO₂'),
        (5.7, 4.5, 'ZnO/Al₂O₃'),
    ]
    for x, h, lbl in high_stacks:
        ax.bar(x, h, width=0.55, color='#EF5350', alpha=0.85, zorder=2)
        ax.text(x, 0.15, lbl, ha='center', va='bottom',
                fontsize=8, color='white', rotation=90, zorder=3)
    ax.text(4.5, 5.1, 'High-CBO (outside window)',
            ha='center', fontsize=10, color='#B71C1C', fontweight='bold')

    low_stacks = [
        (0.60, 1.18, 'In₂O₃/κ-Ga₂O₃', '#1565C0'),
        (1.55, 0.55, 'ZnO/MgZnO ★',    '#0D47A1'),
    ]
    for x, h, lbl, col in low_stacks:
        ax.bar(x, h, width=0.50, color=col, alpha=0.92, zorder=2)
        ax.text(x, 0.1, lbl, ha='center', va='bottom',
                fontsize=8, color='white', rotation=90, zorder=3, fontweight='bold')

    ax.set_xlim(0, 7); ax.set_ylim(0, 5.9)
    ax.set_xlabel('V_res (V)', fontsize=12)
    ax.set_ylabel('Resonance voltage (V)', fontsize=12)
    ax.set_yticks([1, 2, 3, 4, 5])
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.tick_params(labelsize=11)
    plt.tight_layout()
    plt.savefig(out_path, dpi=180, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"  chart → {out_path}")


# ──────────────────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────────────────

def main():
    chart_path = os.path.join(HERE, 'design_rule_chart.png')
    make_bar_chart(chart_path)

    prs = Presentation()
    prs.slide_width  = I(SW)
    prs.slide_height = I(SH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # ── TITLE BAR ──────────────────────────────────────────
    TB_H = 2.40
    add_rect(slide, 0, 0, SW, TB_H, fill=NAVY)

    add_tb(slide, 0.35, 0.12, SW-0.7, 1.0,
           [{'text': 'Material Stack Discovery for BEOL-Compatible Quantum Electronic Nose',
             'size': 50, 'bold': True, 'color': WHITE}],
           ml=0.3, mr=0.2, mt=0.08, mb=0.04)

    add_tb(slide, 0.35, 1.18, SW-0.7, 0.60,
           [{'text': 'Abhineet Agarwal  ·  Swaroop Ganguly',
             'size': 27, 'color': LTBLUE}],
           ml=0.3, mr=0.2, mt=0.04, mb=0.02)

    add_tb(slide, 0.35, 1.75, SW-0.7, 0.50,
           [{'text': ('Department of Electrical Engineering, '
                      'Indian Institute of Technology Bombay, Mumbai, India   |   '
                      'swaroop.ganguly@iitb.ac.in   |   EMRS Spring Meeting 2026'),
             'size': 17, 'italic': True, 'color': WHITE}],
           ml=0.3, mr=0.2, mt=0.02, mb=0.02)

    # ── LAYOUT ─────────────────────────────────────────────
    BY   = TB_H + 0.20   # body top
    BH   = SH - BY - 0.20
    M    = 0.30          # outer margin
    GAP  = 0.28          # column gap
    BGAP = 0.18          # gap between boxes

    C1X, C1W = M,               9.70
    C2X, C2W = C1X+C1W+GAP,   12.30
    C3X       = C2X+C2W+GAP
    C3W       = SW - C3X - M   # ≈ 9.73

    # ════════════════════════════════════════════════════════
    # COLUMN 1
    # ════════════════════════════════════════════════════════
    cy = BY

    # Box 1: Motivation
    B1H = 10.2
    bt = section_box(slide, C1X, cy, C1W, B1H, 'Motivation: Room-Temperature IETS',
                     NAVY, title_sz=20)
    add_tb(slide, C1X+0.15, bt+0.12, C1W-0.30, 2.4,
           [{'text': ('IETS identifies molecules via vibrational fingerprints — peaks in '
                      'd²I/dV² at bias V = ℏω/e for each vibrational mode.'),
             'size': 15},
            {'text': ''},
            {'text': ('Fundamental problem.  Conventional MIM junctions require T < 10 K: '
                      'thermal broadening (~k_BT) washes out vibrational peaks [1].'),
             'size': 15}])

    # callout
    add_rect(slide, C1X+0.15, bt+2.65, C1W-0.30, 0.82,
             fill=LTORANGE, border=ORANGE, bw=1.5)
    add_tb(slide, C1X+0.15, bt+2.65, C1W-0.30, 0.82,
           [{'text': ('RTD solution [2]: the quantum-well resonance acts as an energy filter '
                      'narrower than k_BT, restoring d²I/dV² peaks at 300 K.'),
             'size': 15, 'bold': True, 'color': ORANGE}],
           ml=0.18, mr=0.18, mt=0.12, mb=0.06)

    add_tb(slide, C1X+0.15, bt+3.60, C1W-0.30, 4.2,
           [{'text': ('Remaining gap.  Patil\'s proof-of-concept [2] used GaAs/AlAs grown '
                      'by MBE at >700 °C — incompatible with CMOS back-end-of-line (BEOL).'),
             'size': 15},
            {'text': ''},
            {'text': 'This work:', 'size': 15, 'bold': True, 'color': NAVY},
            {'text': ('Systematic computational screening of 12 oxide heterostructures '
                      '(all ALD/sputtering, <400 °C) to identify BEOL-compatible stacks '
                      'for room-temperature RTD-IETS molecular sensing.'),
             'size': 15}])
    cy += B1H + BGAP

    # Box 2: Selection Criteria
    B2H = 10.2
    bt = section_box(slide, C1X, cy, C1W, B2H, 'Material Selection Criteria',
                     NAVY, title_sz=20)

    criteria = [
        ('1.', 'BEOL compatibility:',    'process temperature < 400 °C'),
        ('2.', 'Band offset:',           'CBO ≈ 0.5 eV — aligns V_res with molecular vibrations (50–400 meV)'),
        ('3.', 'Effective mass:',        'enables quantum confinement in practical barrier/well widths (1–3 nm)'),
        ('4.', 'Mature deposition:',     'ALD or sputtering with demonstrated thickness uniformity'),
        ('5.', 'Integration readiness:', 'prior use in semiconductor industry processes'),
    ]
    ry = bt + 0.18
    for num, lbl, desc in criteria:
        add_tb(slide, C1X+0.15, ry, C1W-0.30, 0.90,
               [{'text': f'{num}  {lbl} {desc}', 'size': 14}])
        ry += 0.90

    add_tb(slide, C1X+0.15, bt+5.7, C1W-0.30, 1.6,
           [{'text': 'Well materials screened:', 'size': 14, 'bold': True, 'color': NAVY},
            {'text': ('In₂O₃ (m*=0.30, LO 70 meV),  IGZO (m*=0.34, LO 60 meV),  '
                      'ZnO (m*=0.28, LO 72 meV),  SnO₂ (m*=0.25, LO 75 meV)'),
             'size': 14}])
    add_tb(slide, C1X+0.15, bt+7.5, C1W-0.30, 1.6,
           [{'text': 'Barrier materials screened:', 'size': 14, 'bold': True, 'color': NAVY},
            {'text': ('Al₂O₃ (CBO 2.8–3.0 eV),  HfO₂ (CBO 2.0–2.3 eV),  '
                      'Ga₂O₃ polymorphs (CBO 0.07–0.61 eV),  '
                      'MgZnO (CBO = 1.57·x_Mg, tunable)'),
             'size': 14}])
    cy += B2H + BGAP

    # Box 3: Design Rule (takes remaining column height)
    B3H = BH - (cy - BY)
    bt = section_box(slide, C1X, cy, C1W, B3H,
                     'Key Insight: The CBO ≈ 0.5 eV Design Rule',
                     NAVY, title_sz=19)
    chart_h = px_img(slide, chart_path, C1X+0.18, bt+0.15, C1W-0.36)
    add_tb(slide, C1X+0.15, bt+chart_h+0.30, C1W-0.30, 2.5,
           [{'text': ('High-CBO oxides resonate at 3–5 V, far outside the molecular '
                      'vibration window (50–400 mV). Only low-CBO stacks (CBO ≈ 0.5 eV) '
                      'bring V_res into the sensing range.'),
             'size': 14}])

    # ════════════════════════════════════════════════════════
    # COLUMN 2
    # ════════════════════════════════════════════════════════
    cy = BY

    # Box 4: Material Survey Table
    B4H = 15.2
    bt = section_box(slide, C2X, cy, C2W, B4H,
                     'Material Survey: Simulation Results (300 K, Symmetric RTD)',
                     NAVY, title_sz=20)
    add_tb(slide, C2X+0.15, bt+0.12, C2W-0.30, 0.75,
           [{'text': ('All devices: 10 nm n⁺ contacts / barrier / 3 nm well / barrier / '
                      '10 nm n⁺ contacts. High-barrier stacks show artifact PVCR only.'),
             'size': 13}])
    make_survey_table(slide, C2X+0.15, bt+0.95, C2W-0.30, 7.0)

    # legend
    add_tb(slide, C2X+0.15, bt+8.08, C2W-0.30, 0.45,
           [{'text': '† Artifact — resonance outside molecular vibration window; not physical NDR.',
             'size': 11, 'italic': True}])
    lx = C2X + 0.15
    for col, lbl in [(ROWRED,'Outside window'), (ROWYELLOW,'CBO too low'),
                     (LTBLUE,'Good candidate'), (ROWBLUE,'Selected')]:
        add_rect(slide, lx, bt+8.65, 0.26, 0.22, fill=col, border=DKGRAY, bw=0.5)
        add_tb(slide, lx+0.30, bt+8.60, 1.75, 0.32,
               [{'text': lbl, 'size': 11}], ml=0.02, mr=0.02, mt=0.02, mb=0.02)
        lx += 2.05

    add_rect(slide, C2X+0.15, bt+9.1, C2W-0.30, 0.85,
             fill=LTORANGE, border=ORANGE, bw=1.5)
    add_tb(slide, C2X+0.15, bt+9.1, C2W-0.30, 0.85,
           [{'text': ('ZnO/Mg₀.₃Zn₀.₇O has the highest PVCR (6.3, comparable to III-V RTDs), '
                      'demonstrated experimental NDR [3,4], and lowest process temperature (<350 °C).'),
             'size': 14, 'bold': True, 'color': ORANGE}],
           ml=0.18, mr=0.18, mt=0.12, mb=0.06)
    cy += B4H + BGAP

    # Box 5: Candidate Stacks
    B5H = 8.8
    bt = section_box(slide, C2X, cy, C2W, B5H,
                     'Candidate Stacks: IETS Viability Criteria',
                     NAVY, title_sz=20)
    make_candidate_table(slide, C2X+0.15, bt+0.12, C2W-0.30, 4.5)
    add_tb(slide, C2X+0.15, bt+4.75, C2W-0.30, 0.55,
           [{'text': '‡ Ref. [3] uses x = 0.2; this work uses x = 0.3 (ΔEc = 1.57·x_Mg [12]).',
             'size': 11, 'italic': True}])
    add_tb(slide, C2X+0.15, bt+5.45, C2W-0.30, 1.5,
           [{'text': ('Only ZnO/Mg₀.₃Zn₀.₇O satisfies all three IETS criteria simultaneously: '
                      'CBO in window, experimental NDR, and BEOL-compatible deposition.'),
             'size': 14}])
    cy += B5H + BGAP

    # Box 6: Simulation Methodology (fills remaining col 2 height)
    B6H = BH - (cy - BY)
    bt = section_box(slide, C2X, cy, C2W, B6H, 'Simulation Methodology',
                     NAVY, title_sz=20)
    sx, sw = C2X+0.15, C2W-0.30

    def method_subbox(y, title, body_text, bh):
        add_rect(slide, sx, bt+y, sw, bh, fill=LTGRAY, border=MIDBLUE, bw=1.0)
        add_rect(slide, sx, bt+y, sw, 0.36, fill=MIDBLUE)
        add_tb(slide, sx+0.1, bt+y, sw-0.2, 0.36,
               [{'text': title, 'size': 16, 'bold': True, 'color': WHITE}],
               ml=0.12, mr=0.06, mt=0.04, mb=0.02)
        add_tb(slide, sx+0.12, bt+y+0.38, sw-0.24, bh-0.44,
               [{'text': body_text, 'size': 14}])

    method_subbox(0.15,
        '1.  NEGF + Rank-1 SCBA',
        ("Quasi-3D Non-Equilibrium Green's Function with self-consistent Born approximation. "
         "For a Gaussian molecular form factor (σ_mol = 3 Å ≪ L⊥), the coupling matrix is "
         "rank-1: M = D₀·χ(z)|u⟩⟨u|  (O(10⁻⁶) correction). Full 3D SCBA collapses to a "
         "single 1D projected Dyson equation, exact for any transverse grid size."),
        3.40)

    method_subbox(3.72,
        '2.  Anderson Mixing (DIIS) [6]',
        ("Fixed-point iteration diverges near NDR. History depth M = 8; solve "
         "min‖Σαᵢrᵢ‖², Σαᵢ = 1 at each step. "
         "10–60 iterations at all 201 bias points; "
         "current conservation |I_L+I_R|/max|I| < 0.2 %."),
        3.10)

    method_subbox(6.99,
        '3.  Numerical d²I/dV²',
        ("d²I/dV² computed by diff(diff(I))/dV² on the converged SCBA I(V) — "
         "exactly the procedure of Patil [2]. Self-consistent SCBA produces "
         "inherently smooth I(V); no smoothing applied."),
        2.65)

    # flowchart
    fc_path = os.path.join(HERE, 'fig3_flowchart.png')
    if os.path.exists(fc_path):
        fy = 9.82
        fh = px_img(slide, fc_path, sx+0.1, bt+fy, sw-0.2)
        add_tb(slide, sx, bt+fy+fh+0.05, sw, 0.35,
               [{'text': 'Fig. 1. Rank-1 projected NEGF–SCBA convergence loop.',
                 'size': 11, 'italic': True}], ml=0.05, mr=0.05, mt=0.02, mb=0.02)

    # ════════════════════════════════════════════════════════
    # COLUMN 3
    # ════════════════════════════════════════════════════════
    cy = BY

    # Box 7: Selected Stack
    B7H = 9.0
    bt = section_box(slide, C3X, cy, C3W, B7H,
                     'Selected Stack: ZnO/Mg₀.₃Zn₀.₇O RTD',
                     GREEN, title_sz=19)
    add_tb(slide, C3X+0.15, bt+0.12, C3W-0.30, 2.1,
           [{'text': ('Symmetric 2 nm MgZnO / 3 nm ZnO / 2 nm MgZnO, '
                      '10 nm n⁺-ZnO contacts (N_D = 10²⁵ m⁻³).  '
                      'ΔEc = 0.47 eV,  m* = 0.28 m₀,  ℏω_LO = 72 meV.  '
                      'Single bound state E₁ = 149 meV; NDR at V_res ≈ 552 mV, '
                      'peak ~72 nA, PVR ~6 (SCBA, 300 K).'),
             'size': 14}])
    bd_path = os.path.join(HERE, 'fig2_band_diagram.png')
    if os.path.exists(bd_path):
        bh_ = px_img(slide, bd_path, C3X+0.15, bt+2.35, C3W-0.30)
        add_tb(slide, C3X+0.15, bt+2.35+bh_+0.05, C3W-0.30, 0.6,
               [{'text': 'Fig. 2. Band diagram: (i) V = 0; (ii) resonance at 552 mV. Orange Gaussian χ(z): molecular coupling region.',
                 'size': 11, 'italic': True, 'align': PP_ALIGN.CENTER}],
               ml=0.05, mr=0.05, mt=0.02, mb=0.02)
    cy += B7H + BGAP

    # Box 8: Hero Result
    B8H = 13.2
    bt = section_box(slide, C3X, cy, C3W, B8H,
                     'Key Result: Molecular Discrimination at 300 K',
                     ORANGE, title_sz=19)
    fw = (C3W - 0.42) / 2.0
    h4 = h5 = 4.0
    iv_path  = os.path.join(HERE, 'fig4_IV.png')
    d2_path  = os.path.join(HERE, 'fig5_d2IdV2.png')
    if os.path.exists(iv_path):
        h4 = px_img(slide, iv_path,  C3X+0.15, bt+0.12, fw)
    if os.path.exists(d2_path):
        h5 = px_img(slide, d2_path, C3X+0.15+fw+0.12, bt+0.12, fw)
    fig_bot = bt + max(h4, h5) + 0.12
    add_tb(slide, C3X+0.15, fig_bot, fw, 0.65,
           [{'text': '(a) I–V. NDR at ≈ 552 mV, PVR ~6. Molecule-specific ΔI in pre-NDR region.',
             'size': 10, 'italic': True, 'align': PP_ALIGN.CENTER}],
           ml=0.03, mr=0.03, mt=0.02, mb=0.02)
    add_tb(slide, C3X+0.15+fw+0.12, fig_bot, fw, 0.65,
           [{'text': '(b) d²I/dV². Peaks at 100 mV (Mol A) and 180 mV (Mol B) at V = ℏω/e.',
             'size': 10, 'italic': True, 'align': PP_ALIGN.CENTER}],
           ml=0.03, mr=0.03, mt=0.02, mb=0.02)

    dd_path = os.path.join(HERE, 'fig7_deltaD.png')
    dd_y = fig_bot + 0.72
    if os.path.exists(dd_path):
        hd = px_img(slide, dd_path, C3X+0.15, dd_y, C3W-0.30)
        add_tb(slide, C3X+0.15, dd_y+hd+0.05, C3W-0.30, 0.75,
               [{'text': ('Fig. 3. ΔD = d²I_mol/dV² − d²I_BL/dV². '
                          'Each molecule gives a distinct peak pattern. '
                          'Mol AB ≈ superposition of A + B (rank-1 linearity).'),
                 'size': 10, 'italic': True, 'align': PP_ALIGN.CENTER}],
               ml=0.05, mr=0.05, mt=0.02, mb=0.02)

    call_y = bt + B8H - 0.95
    add_rect(slide, C3X+0.15, call_y, C3W-0.30, 0.82,
             fill=LTORANGE, border=ORANGE, bw=1.5)
    add_tb(slide, C3X+0.15, call_y, C3W-0.30, 0.82,
           [{'text': 'Three synthetic odorants discriminated at 300 K — no cryogenic cooling.',
             'size': 15, 'bold': True, 'color': ORANGE, 'align': PP_ALIGN.CENTER}],
           ml=0.18, mr=0.18, mt=0.14, mb=0.06)
    cy += B8H + BGAP

    # Box 9: Temperature
    B9H = 6.8
    bt = section_box(slide, C3X, cy, C3W, B9H,
                     'Temperature Dependence — Mol A (ℏω = 100 meV)',
                     MIDBLUE, title_sz=18)
    td_path = os.path.join(HERE, 'fig8_temp_dependence.png')
    if os.path.exists(td_path):
        th = px_img(slide, td_path, C3X+0.15, bt+0.12, C3W-0.30)
        add_tb(slide, C3X+0.15, bt+0.12+th+0.05, C3W-0.30, 1.0,
               [{'text': ('Fig. 4. 10–300 K. (a) Peak current: 73 nA (300 K) → 110 nA (10 K). '
                          '(b) IETS peak sharpens at low T but remains resolvable at 300 K — '
                          'RTD energy filtering preserves features beyond the 3k_BT broadening limit [2].'),
                 'size': 10, 'italic': True, 'align': PP_ALIGN.CENTER}],
               ml=0.05, mr=0.05, mt=0.02, mb=0.02)
    cy += B9H + BGAP

    # Box 10: Fabrication
    B10H = 6.8
    bt = section_box(slide, C3X, cy, C3W, B10H, 'Fabrication: ALD Deposition Routes',
                     NAVY, title_sz=19)
    make_ald_table(slide, C3X+0.15, bt+0.12, C3W-0.30, 3.2)
    add_tb(slide, C3X+0.15, bt+3.45, C3W-0.30, 2.8,
           [{'text': ('All layers depositable in a single ALD tool run below 350 °C — '
                      'within the BEOL thermal budget.  MgZnO composition is controlled by '
                      'the Mg(Cp)₂:DEZ pulse ratio (ΔEc = 1.57·x_Mg).  '
                      'Wurtzite phase-stable at x_Mg ≤ 0.33, avoiding phase segregation.  '
                      'Proposed device: ALD stack on CMOS backend test structure.'),
             'size': 13}])
    cy += B10H + BGAP

    # Box 11: Conclusions
    remaining = BH - (cy - BY)
    B11H = remaining * 0.43
    bt = section_box(slide, C3X, cy, C3W, B11H, 'Conclusions',
                     NAVY, title_sz=20)
    bullets = [
        'High-CBO oxides (CBO > 2 eV) resonate at 3–5 V — outside molecular vibration range.',
        'Design rule: CBO ≈ 0.5 eV aligns V_res with the 50–400 meV fingerprint band.',
        'ZnO/Mg₀.₃Zn₀.₇O: only oxide stack with experimental NDR + BEOL ALD + highest PVCR (6.3).',
        'NEGF–SCBA with Anderson mixing discriminates three odorants at 300 K.',
        'Full ALD fabrication route identified; next step: RTD on CMOS test structure.',
    ]
    by_ = bt + 0.18
    for b in bullets:
        add_tb(slide, C3X+0.18, by_, C3W-0.30, 0.80,
               [{'text': '▶  ' + b, 'size': 13}])
        by_ += 0.80
    cy += B11H + BGAP

    # Box 12: References
    B12H = BH - (cy - BY)
    bt = section_box(slide, C3X, cy, C3W, B12H, 'References',
                     NAVY, title_sz=19)
    refs = [
        '[1] R.C. Jaklevic & J. Lambe, PRL 17, 1139 (1966).',
        '[2] A. Patil et al., Sci. Rep. 8, 128 (2018).',
        '[3] S. Krishnamoorthy et al., SSE 46, 1633 (2002).',
        '[4] I.-C. Cheng et al., IEEE EDL 36, 1330 (2015).',
        '[5] S. Datta, Quantum Transport, Cambridge (2005).',
        '[6] D.G. Anderson, J. ACM 12, 547 (1965).',
        '[7] N. Pandey et al., Sci. Rep. 11, 11389 (2021).',
        '[8] T.C.L.G. Sollner et al., APL 43, 588 (1983).',
        '[9] C. Bayram et al., APL 96, 042103 (2010).',
        '[10] Y. Kuang et al., ACS AEM 3, 795 (2021).',
        '[11] J. Zhu et al., Nanotechnology 35 (2024).',
        '[12] H. Yin et al., Sci. Rep. 7, 41567 (2017).',
    ]
    half = len(refs) // 2
    rw   = (C3W - 0.36) / 2.0
    for i, ref in enumerate(refs[:half]):
        add_tb(slide, C3X+0.15, bt+0.12+i*0.50, rw, 0.50,
               [{'text': ref, 'size': 10}], ml=0.05, mr=0.02, mt=0.02, mb=0.02)
    for i, ref in enumerate(refs[half:]):
        add_tb(slide, C3X+0.15+rw+0.06, bt+0.12+i*0.50, rw, 0.50,
               [{'text': ref, 'size': 10}], ml=0.05, mr=0.02, mt=0.02, mb=0.02)

    # ── Save ───────────────────────────────────────────────
    out = os.path.join(HERE, 'emrs_poster.pptx')
    prs.save(out)
    print(f"\nSaved: {out}")
    print(f"Slide: {SW}\" × {SH}\" (A0 portrait)")
    print("Open in PowerPoint/Keynote → File → Print → A0 paper.")


if __name__ == '__main__':
    main()
